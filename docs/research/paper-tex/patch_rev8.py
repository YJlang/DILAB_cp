"""rev7 → rev8. 교수님 피드백: 실험 플랫폼(DILAB) 소개 슬라이드 추가.

Performance Evaluation 앞에 새 슬라이드 한 장을 끼워, 실험 코퍼스가 나온
운영 서비스(DILAB)를 구조도 + 대시보드 화면으로 먼저 보여준다. 두 그림은
연구실 공개 페이지(inc.cau.ac.kr)에 이미 게시된 원본이라 공개 범위 문제가 없다.

삽입 방식: 덱의 장식(파란 탭·회색 바·빨간 쪽번호 블록)은 레이아웃이 아니라
슬라이드 위 도형이므로, 새 레이아웃 슬라이드를 만들면 맨몸이 된다. 그래서
이웃 슬라이드에서 장식 도형만 XML 복사해 온다 — 장식은 rel(그림 참조)이
없어서 통짜 spTree 복사와 달리 깨질 링크가 생기지 않는다.
"""
import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev7.pptx"
DST = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev8.pptx"
ARCH = ROOT / "docs/figures/fig1-architecture.png"
DASH = ROOT / "docs/figures/fig2-product.png"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6D, 0x6D, 0x6D)
BLUE = RGBColor(0x33, 0x75, 0xBE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
BODY_L, BODY_W, BODY_T = 0.78, 11.29, 1.52
GREY_BAR_TOP = 7.077
IMG_GAP = 0.30                 # 두 그림 사이 간격(in)

INSERT_AT = 7                  # 기존 Performance Evaluation 자리 (0-based)


def IN(v):
    return Emu(v).inches if v is not None else 0.0


def _style(run, size, bold, color):
    f = run.font
    f.size, f.bold, f.color.rgb, f.name = Pt(size), bold, color, FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "en-US")
    rPr.set("altLang", "en-US")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def bullets(slide, items, top, size=16):
    tf = textbox(slide, BODY_L, top, BODY_W, 0.40 * len(items) + 0.30)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(9)
        p.line_spacing = 1.25
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(0.28 * 914400)))
        pPr.set("indent", str(int(-0.28 * 914400)))
        _style(p.add_run(), size, False, INK)
        p.runs[0].text = "·   " + it
    return top + 0.40 * len(items) + 0.30


def drop(sh):
    sh._element.getparent().remove(sh._element)


def is_decoration(sh):
    return sh.shape_type in (1, 9) and sh.name.startswith(("직사각형", "직선"))


def set_pagenum(slide, n):
    for sh in slide.shapes:
        if (sh.shape_type == 1 and abs(IN(sh.left) - 12.66) < .06
                and abs(IN(sh.top) - 7.08) < .06):
            tf = sh.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _style(p.add_run(), 12, True, WHITE)
            p.runs[0].text = str(n)
            return True
    return False


def caption(slide, text, left, top, width):
    tf = textbox(slide, left, top, width, 0.28)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _style(p.add_run(), 11, False, GREY)
    p.runs[0].text = text


def build():
    prs = Presentation(str(SRC))
    assert len(prs.slides._sldIdLst) == 16, "rev7 은 16장이어야 한다"
    donor = prs.slides[INSERT_AT]          # Performance Evaluation — 장식·제목 스타일의 원본

    # ── 새 슬라이드: 레이아웃으로 만들고 장식만 이웃에서 복사 ──────────
    sl = prs.slides.add_slide(donor.slide_layout)
    for sh in list(sl.shapes):             # 레이아웃이 심은 placeholder 제거
        drop(sh)
    for sh in donor.shapes:
        if is_decoration(sh):
            sl.shapes._spTree.append(copy.deepcopy(sh._element))

    # 제목은 이웃 슬라이드의 실제 문구를 그대로 쓴다. 사용자가 절 번호를
    # 직접 뗀 이력이 있어(rev6 참조) 하드코딩하면 어긋난다.
    title_text = next(sh.text_frame.text.strip() for sh in donor.shapes
                      if sh.has_text_frame and not is_decoration(sh)
                      and abs(IN(sh.top) - 0.246) < .06)
    tf = textbox(sl, BODY_L, 0.246, BODY_W, 0.60)
    _style(tf.paragraphs[0].add_run(), 24, True, INK)
    tf.paragraphs[0].runs[0].text = title_text

    tf = textbox(sl, BODY_L, 0.990, BODY_W, 0.42)
    _style(tf.paragraphs[0].add_run(), 17, True, BLUE)
    tf.paragraphs[0].runs[0].text = (
        "Testbed — DILAB, a product-review analysis service in production")

    y = bullets(sl, [
        "Consumer and expert reviews are embedded into a vector store "
        "and served as evidence-cited reports",
        "The 1,347-chunk corpus and 43 queries in our experiments "
        "come from this service",
    ], BODY_T)

    # ── 그림 두 장: 같은 높이로 나란히, 본문 폭을 꽉 채운다 ───────────
    ar_a = Image.open(ARCH).size[0] / Image.open(ARCH).size[1]
    ar_d = Image.open(DASH).size[0] / Image.open(DASH).size[1]
    h = (BODY_W - IMG_GAP) / (ar_a + ar_d)
    w_a, w_d = h * ar_a, h * ar_d
    top = y + 0.30
    cap_top = top + h + 0.08
    assert cap_top + 0.28 <= GREY_BAR_TOP, "캡션이 회색 바를 파고든다"

    sl.shapes.add_picture(str(ARCH), Inches(BODY_L), Inches(top),
                          Inches(w_a), Inches(h))
    sl.shapes.add_picture(str(DASH), Inches(BODY_L + w_a + IMG_GAP), Inches(top),
                          Inches(w_d), Inches(h))
    caption(sl, "System architecture & analysis pipeline", BODY_L, cap_top, w_a)
    caption(sl, "Product report dashboard", BODY_L + w_a + IMG_GAP, cap_top, w_d)

    # ── 맨 끝에 태어난 슬라이드를 제자리로 옮기고 쪽번호 재부여 ────────
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[-1]
    lst.remove(el)
    lst.insert(INSERT_AT, el)

    for i, sl2 in enumerate(prs.slides):
        if i == 0:
            continue                       # 표지는 쪽번호가 없다
        set_pagenum(sl2, i)

    assert len(prs.slides._sldIdLst) == 17
    prs.save(str(DST))
    return DST, w_a, w_d, h


if __name__ == "__main__":
    out, w_a, w_d, h = build()
    print(f"저장: {out}")
    print(f"그림: 구조도 {w_a:.2f}×{h:.2f}in · 대시보드 {w_d:.2f}×{h:.2f}in")
