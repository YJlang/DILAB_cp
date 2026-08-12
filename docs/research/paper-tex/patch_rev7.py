"""rev6 → rev7. 랩장님 최종 피드백 2건.

  (1) Conclusion 을 한 페이지로
  (2) References 폰트 키우기

(1)에는 함정이 있다. Conclusion 두 장을 그대로 합치면 꼭지가 8개가 되어
바로 앞에 받은 "꼭지 2~4개" 규칙을 깬다. 그래서 꼭지는 핵심 4개만 남기고,
한계·향후과제는 꼭지가 아닌 작은 각주 한 문단으로 내린다. 논문이 리뷰어에게
평가받은 지점이 "한계를 분명히 밝힌 것"이라 통째로 버릴 수는 없다.

(2)는 9건이 한 장에 들어가야 하므로 들어가는 최대 크기를 계산해서 정한다.
"""
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev6.pptx"
DST = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev7.pptx"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6D, 0x6D, 0x6D)
BLUE = RGBColor(0x33, 0x75, 0xBE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
BODY_L, BODY_W, BODY_T = 0.78, 11.29, 1.52
REF_BOTTOM = 6.78


def IN(v):
    return Emu(v).inches if v is not None else 0.0


def _style(run, size, bold, color):
    f = run.font
    f.size, f.bold, f.color.rgb, f.name = Pt(size), bold, color, FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "en-US")
    rPr.set("altLang", "en-US")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def bullets(slide, items, top, size=16):
    tf = textbox(slide, BODY_L, top, BODY_W, 0.40 * len(items) + 0.30)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(9)
        p.line_spacing = 1.25
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(0.28 * 914400)))
        pPr.set("indent", str(int(-0.28 * 914400)))
        _style(p.add_run(), size, False, INK)
        p.runs[0].text = "·   " + it
    return top + 0.40 * len(items) + 0.30


def subhead(slide, text, top):
    tf = textbox(slide, BODY_L, top, BODY_W, 0.40)
    _style(tf.paragraphs[0].add_run(), 17, True, BLUE)
    tf.paragraphs[0].runs[0].text = text
    return top + 0.50


def drop(sh):
    sh._element.getparent().remove(sh._element)


def content_shapes(slide):
    """장식이 아닌 도형만."""
    return [sh for sh in slide.shapes
            if not (sh.shape_type in (1, 9) and sh.name.startswith(("직사각형", "직선")))]


def set_pagenum(slide, n):
    for sh in slide.shapes:
        if (sh.shape_type == 1 and abs(IN(sh.left) - 12.66) < .06
                and abs(IN(sh.top) - 7.08) < .06):
            tf = sh.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _style(p.add_run(), 12, True, WHITE)
            p.runs[0].text = str(n)
            return True
    return False


def drop_slide(prs, idx):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    prs.part.drop_rel(ids[idx].get(qn("r:id")))
    lst.remove(ids[idx])


REFS = [
    ('[1] P. Lewis et al., “Retrieval-augmented generation for knowledge-intensive '
     'NLP tasks,” Advances in Neural Information Processing Systems, vol. 33, '
     'pp. 9459–9474, 2020.'),
    ('[2] Y. Gao et al., “Retrieval-augmented generation for large language models: '
     'A survey,” arXiv preprint arXiv:2312.10997, 2023.'),
    ('[3] Oracle, “Oracle AI Vector Search User’s Guide,” '
     'Oracle AI Database 26ai documentation, 2026.'),
    ('[4] L. Wang, N. Yang, X. Huang, L. Yang, R. Majumder, and F. Wei, '
     '“Multilingual E5 text embeddings: A technical report,” '
     'arXiv preprint arXiv:2402.05672, 2024.'),
    ('[5] pgvector, “pgvector: Open-source vector similarity search for Postgres,” '
     'GitHub repository, 2026.'),
    ('[6] J. Chen, S. Xiao, P. Zhang, K. Luo, D. Lian, and Z. Liu, '
     '“M3-Embedding: Multi-linguality, multi-functionality, multi-granularity text '
     'embeddings through self-knowledge distillation,” Findings of the Association '
     'for Computational Linguistics: ACL 2024, pp. 2318–2335, 2024.'),
    ('[7] Y. A. Malkov and D. A. Yashunin, “Efficient and robust approximate nearest '
     'neighbor search using hierarchical navigable small world graphs,” IEEE '
     'Transactions on Pattern Analysis and Machine Intelligence, vol. 42, no. 4, '
     'pp. 824–836, 2020.'),
    ('[8] L. Zheng et al., “Judging LLM-as-a-judge with MT-Bench and Chatbot Arena,” '
     'Advances in Neural Information Processing Systems, vol. 36, '
     'Datasets and Benchmarks Track, 2023.'),
    ('[9] D. Lakens, “Equivalence tests: A practical primer for t tests, correlations, '
     'and meta-analyses,” Social Psychological and Personality Science, vol. 8, '
     'no. 4, pp. 355–362, 2017.'),
]

HANG = 0.42
GAP_PT = 9


def refs_height(size):
    """들여쓰기·줄바꿈을 감안한 참고문헌 블록 높이(in).

    라틴 문자 평균 폭 계수는 0.56이다. 0.50으로 잡으면 줄 수를 과소평가해
    한 단계 큰 글자를 고르고, 마지막 항목이 회색바를 파고든다.
    """
    usable = BODY_W - HANG
    per_char = size * 0.56 / 72
    total = 0.0
    for r in REFS:
        lines = max(1, math.ceil(len(r) * per_char / usable))
        total += lines * size * 1.18 / 72
    return (total + (len(REFS) - 1) * GAP_PT / 72) * 1.08   # 줄바꿈 지점 오차 여유


def build():
    prs = Presentation(str(SRC))
    s = prs.slides

    # ── s14 · Conclusion 한 장으로 통합 ────────────────────────────────
    sl = s[13]
    for sh in content_shapes(sl):
        if abs(IN(sh.top) - 0.25) < .12 or abs(IN(sh.top) - 0.99) < .12:
            continue                       # 제목·소제목은 그대로 둔다
        drop(sh)

    y = bullets(sl, [
        "A single bridge condition splits an end-to-end comparison into two single-factor comparisons",
        "No store effect was detected; the entire observed change came from the placement",
    ], BODY_T)
    y = subhead(sl, "A practical decision rule", y + 0.22)
    y = bullets(sl, [
        "If governance of the stored data is the only goal, replacing the store alone is enough",
        "If the query path must stay inside, confirm the model size limit before migrating",
    ], y)

    # 한계·향후과제는 꼭지가 아닌 각주로. 꼭지 수를 4개로 지키면서 내용은 남긴다.
    # 각 줄을 한 줄에 들어가는 길이로 끊는다. 줄바꿈이 없으면 단어가 잘릴
    # 자리도 생기지 않는다 (13pt · 폭 11.29in 기준 약 125자).
    tf = textbox(sl, BODY_L, y + 0.40, BODY_W, 1.10)
    for i, (lead, rest) in enumerate([
        ("Limitations — ",
         "one Korean corpus of 1,347 chunks, relevance scored without human assessors, "
         "an entry-level instance."),
        ("", "Placement and model size are coupled by construction, and the store "
             "result does not establish equivalence."),
        ("Future work — ",
         "a fourth condition running the smaller model outside the database."),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(6)
        p.line_spacing = 1.25
        if lead:
            _style(p.add_run(), 13, True, GREY)
            p.runs[-1].text = lead
        _style(p.add_run(), 13, False, GREY)
        p.runs[-1].text = rest

    drop_slide(prs, 14)                    # 옛 Conclusion 두 번째 장 제거

    # ── References 폰트 키우기 ────────────────────────────────────────
    ref_i = next(i for i, sl2 in enumerate(prs.slides)
                 if any(sh.has_text_frame and sh.text_frame.text.strip() == "References"
                        for sh in sh_list(sl2)))
    sl = prs.slides[ref_i]
    # 이 장만 파란 소제목이 없다. 본문 시작을 올려 빈 띠를 없애고 글자를 더 키운다.
    ref_top = 1.15
    # 어림식은 후보를 좁히는 용도이고, 최종 크기는 렌더 실측으로 확정했다.
    # 13pt 는 1.09in 이 남고 15pt 는 회색바를 파고들어 14pt 로 고정한다.
    size = min(14, max(sz for sz in (13, 14, 15, 16)
                       if refs_height(sz) <= (REF_BOTTOM - ref_top) * 1.10))
    for sh in content_shapes(sl):
        if sh.has_text_frame and sh.text_frame.text.lstrip().startswith("[1]"):
            drop(sh)
    tf = textbox(sl, BODY_L, ref_top, BODY_W, REF_BOTTOM - ref_top)
    for i, r in enumerate(REFS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(GAP_PT)
        p.line_spacing = 1.18
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(HANG * 914400)))
        pPr.set("indent", str(int(-HANG * 914400)))
        _style(p.add_run(), size, False, INK)
        p.runs[0].text = r

    # ── 쪽번호 다시 매기기 (한 장 줄었다) ──────────────────────────────
    for i, sl2 in enumerate(prs.slides):
        if i == 0:
            continue
        set_pagenum(sl2, i)

    prs.save(str(DST))
    return DST, size


def sh_list(slide):
    return list(slide.shapes)


if __name__ == "__main__":
    out, size = build()
    print(f"저장: {out}   References {size}pt")
