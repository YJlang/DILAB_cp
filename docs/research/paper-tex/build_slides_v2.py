"""발표 슬라이드 v2 — 피드백 2건 반영본.

rev1을 열어 템플릿 장식(파란 사각형·회색바·빨간 쪽번호블록)만 남기고
내용 도형을 전부 걷어낸 뒤 새로 앉힌다. 장식을 그대로 물려받아야
서식이 원본과 어긋나지 않는다.

반영한 피드백
  [1-1] 쪽번호를 표지 제외 1부터
  [1-2] 2~5쪽 텍스트 → 도식
  [1-3] 그림 가운데 정렬 통일
  [1-4] 소제목 아래 bullet, 그림은 아래쪽
  [1-5] 실험 설정 표 재디자인 (Embedding model 줄바꿈 제거)
  [1-6] 결과 슬라이드에 시각화 추가
  [2-1] 한글 전부 제거
  [2-2] 마지막 장은 Thank You 만
  [2-3] 단어 중간 끊김 수정 (lang/폰트 지정)
  [2-4] 논문 형태의 절 제목
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev1.pptx"
DST = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev3.pptx"
FIG = ROOT / "docs/figures"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6D, 0x6D, 0x6D)
BLUE = RGBColor(0x33, 0x75, 0xBE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC8, 0xC8, 0xC8)

FONT = "Arial"          # 현장 Windows PC·Mac 양쪽에 있는 폰트만 쓴다
SLIDE_W = 13.333
TITLE_T, SUB_T, BODY_T = 0.30, 0.99, 1.52
FIG_BOTTOM = 6.78       # 회색바(7.08) 위 여유
BODY_L, BODY_W = 0.78, 11.29
FIG_MAX_W = 11.60


# ── 텍스트 유틸 ──────────────────────────────────────────────────────────
def _style(run, size, bold, color, italic=False):
    f = run.font
    f.size, f.bold, f.italic = Pt(size), bold, italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    # lang 을 영어로 박아야 PowerPoint 가 동아시아 줄바꿈 규칙을 적용하지 않는다.
    # 이게 없으면 'hosted' 가 'host / ed' 로 잘린다 (피드백 2-3).
    rPr.set("lang", "en-US")
    rPr.set("altLang", "en-US")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def put(tf, text, size, bold=False, color=INK, italic=False, align=PP_ALIGN.LEFT,
        space_before=0, first=False, hanging=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if hanging is not None:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(hanging * 914400)))
        pPr.set("indent", str(int(-hanging * 914400)))
    _style(p.add_run(), size, bold, color, italic)
    p.runs[0].text = text
    return p


# ── 슬라이드 유틸 ────────────────────────────────────────────────────────
KEEP = {1, 9}                       # AUTO_SHAPE, LINE
DECO = ("직사각형", "직선")          # 템플릿에서 온 장식만 한글 이름을 갖는다


def wipe(slide):
    """장식만 남기고 내용 도형을 걷어낸다.

    AutoShape 도 text_frame 을 가지므로 has_text_frame 으로 거르면 장식까지
    지워진다. 반대로 종류만 보면 이전 판에서 만든 빈 사각형(수치 카드 배경)이
    남는다. 그래서 종류 + 이름을 함께 본다.
    """
    for sh in list(slide.shapes):
        if sh.shape_type in KEEP and sh.name.startswith(DECO):
            continue
        sh._element.getparent().remove(sh._element)
    return sum(1 for sh in slide.shapes)


def _numblock(slide):
    """쪽번호가 들어가는 빨간 블록. 장식 AutoShape 자체가 글자를 품고 있다."""
    for sh in slide.shapes:
        if (sh.shape_type == 1
                and abs(Emu(sh.left).inches - 12.66) < 0.06
                and abs(Emu(sh.top).inches - 7.08) < 0.06):
            return sh
    return None


def page_number(slide, n):
    """블록 안에 직접 쓴다. 위에 새 텍스트를 얹으면 옛 번호와 겹쳐 찍힌다."""
    blk = _numblock(slide)
    if blk is None:                       # 블록이 없으면 새로 얹는다
        _, tf = textbox(slide, 12.66, 7.08, 0.68, 0.42, anchor=MSO_ANCHOR.MIDDLE)
        put(tf, str(n), 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
        return
    tf = blk.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _style(p.add_run(), 12, True, WHITE)
    p.runs[0].text = "" if n is None else str(n)


def head(slide, title, sub=None):
    _, tf = textbox(slide, 0.74, TITLE_T, BODY_W, 0.60)
    put(tf, title, 24, bold=True, color=INK, first=True)
    if sub:
        _, tf2 = textbox(slide, BODY_L, SUB_T, BODY_W, 0.42)
        put(tf2, sub, 17, bold=True, color=BLUE, first=True)


def bullets(slide, items, top=BODY_T, size=16, width=BODY_W):
    """반환: 목록이 끝나는 y 좌표."""
    h = 0.40 * len(items) + 0.30
    _, tf = textbox(slide, BODY_L, top, width, h)
    for i, it in enumerate(items):
        p = put(tf, "·   " + it, size, color=INK, first=(i == 0),
                space_before=0 if i == 0 else 9, hanging=0.28)
        p.line_spacing = 1.25
    return top + h


def figure(slide, png, bottom=FIG_BOTTOM, max_w=FIG_MAX_W, top_limit=None):
    """가로 가운데 정렬 + 아래쪽 기준 배치 (피드백 1-3, 1-4)."""
    iw, ih = Image.open(png).size
    w = max_w
    h = w * ih / iw
    if top_limit is not None and bottom - h < top_limit:
        h = bottom - top_limit
        w = h * iw / ih
    left = (SLIDE_W - w) / 2
    slide.shapes.add_picture(str(png), Inches(left), Inches(bottom - h),
                             Inches(w), Inches(h))
    return bottom - h


def hline(slide, l, t, w, color=LINE, pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                               Inches(w), Pt(pt))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def drop_slide(prs, idx):
    xml = prs.slides._sldIdLst
    ids = list(xml)
    rId = ids[idx].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml.remove(ids[idx])


# ── 본문 ────────────────────────────────────────────────────────────────
def build():
    prs = Presentation(str(SRC))
    drop_slide(prs, 13)          # 등가성 그림 슬라이드는 본문과 합친다 → 16장
    s = prs.slides

    # 1 · 표지 (쪽번호 없음)
    sl = s[0]
    wipe(sl)
    page_number(sl, None)       # 표지 블록에 남아 있을 수 있는 글자를 비운다
    # 템플릿의 파란 밑줄은 한 줄 제목을 전제로 놓여 있어 3줄 제목을 관통한다.
    for sh in sl.shapes:
        if sh.shape_type == 9:          # LINE
            sh.top = Inches(3.56)
    _, tf = textbox(sl, 0.90, 1.52, 11.53, 1.95, anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate([
            "A Factor Decomposition Method for Separating",
            "Store Substitution and Embedding Placement Effects",
            "in Vector Database Migration"]):
        p = put(tf, ln, 27, bold=True, color=INK, align=PP_ALIGN.CENTER, first=(i == 0))
        p.line_spacing = 1.22
    _, tf = textbox(sl, 1.20, 3.98, 10.93, 2.10)
    put(tf, "2026 International Conference on PyGeek", 15, bold=True,
        color=INK, align=PP_ALIGN.CENTER, first=True)
    put(tf, "August 20, 2026   ·   Hongik University, Seoul", 14,
        color=GREY, align=PP_ALIGN.CENTER, space_before=4)
    put(tf, "Junha Yoon · Ohhyeon Gwon · Junhaeng Lee · Hyeji Roh · "
            "Soowan Cho · Suhee Kim · SangSoon Lim", 14, bold=True,
        color=INK, align=PP_ALIGN.CENTER, space_before=26)
    put(tf, "Sungkyul University · Chung-Ang University · Deep Insight Lab", 12,
        color=GREY, align=PP_ALIGN.CENTER, space_before=5)

    # 2 · Contents (p.1)
    sl = s[1]; wipe(sl); page_number(sl, 1)
    head(sl, "Contents")
    rows = ["Introduction", "Proposed Method", "Performance Evaluation",
            "Experiment Result", "Conclusion"]
    y = 1.72
    for i, label in enumerate(rows, 1):
        _, tf = textbox(sl, 1.65, y, 0.75, 0.62, anchor=MSO_ANCHOR.MIDDLE)
        put(tf, str(i), 30, bold=True, color=BLUE, first=True)
        _, tf = textbox(sl, 2.55, y, 8.5, 0.62, anchor=MSO_ANCHOR.MIDDLE)
        put(tf, label, 23, bold=True, color=INK, first=True)
        if i < len(rows):
            hline(sl, 1.65, y + 0.74, 9.4)
        y += 0.94

    # 3 · Introduction — RAG 구성 (p.2)
    sl = s[2]; wipe(sl); page_number(sl, 2)
    head(sl, "1. Introduction",
         "The vector layer is the part that gets replaced after deployment")
    bullets(sl, [
        "A RAG system combines an embedding model, a store, and a language model",
        "A common migration moves it to a converged database that embeds internally",
    ])
    figure(sl, FIG / "slide2-fig-intro-rag.png")

    # 4 · Introduction — 두 요인 (p.3)
    sl = s[3]; wipe(sl); page_number(sl, 3)
    head(sl, "1. Introduction",
         "One migration changes two things at the same time")
    bullets(sl, [
        "The store and the embedding placement move together, not one at a time",
        "A governance-driven migration pulls the model inside, and the model size limit follows",
        "A single before-and-after comparison cannot attribute the change to either factor",
    ])
    figure(sl, FIG / "slide2-fig-intro-problem.png")

    # 5 · Introduction — 오진단 (p.4)
    sl = s[4]; wipe(sl); page_number(sl, 4)
    head(sl, "1. Introduction",
         "Misattributing the cause leads to the wrong remedy")
    bullets(sl, [
        "A decline cannot be assigned to the store or to the model",
        "No observed change may mean no effect, or two effects cancelling each other out",
    ])
    figure(sl, FIG / "slide2-fig-intro-misdiagnosis.png")

    # 6 · Proposed Method — 브리지 조건 (p.5)
    sl = s[5]; wipe(sl); page_number(sl, 5)
    head(sl, "2. Proposed Method",
         "A bridge condition splits one comparison into two")
    bullets(sl, [
        "Load the target store with the vectors the source system has already produced",
        "A to B differs only in the store, B to C only in the placement",
        "No re-embedding is required, so the running service is left untouched",
    ])
    figure(sl, FIG / "_decomposition.png")

    # 7 · Proposed Method — 구조 (p.6)
    sl = s[6]; wipe(sl); page_number(sl, 6)
    head(sl, "2. Proposed Method",
         "Network boundaries drop from two to one and the external model call disappears")
    bullets(sl, [
        "The retrieval rule and its parameters were held identical across the three conditions",
    ])
    figure(sl, FIG / "_architecture.png", top_limit=2.10)

    # 8 · Performance Evaluation (p.7)
    sl = s[7]; wipe(sl); page_number(sl, 7)
    head(sl, "3. Performance Evaluation",
         "Three conditions measured over 1,347 chunks and 43 fixed queries")
    bullets(sl, [
        "Retrieval agreement, evidence relevance on a 0 to 2 scale, and search latency",
        "Confidence intervals from 10,000 bootstrap resamples on an entry-level instance",
    ])
    figure(sl, FIG / "slide2-fig-setup-matrix.png")

    # 9 · Experiment Result — 저장소 요인 (p.8)
    sl = s[8]; wipe(sl); page_number(sl, 8)
    head(sl, "4. Experiment Result",
         "Store factor — no quality change was detected")
    bullets(sl, [
        "An identical top-10 set was returned for 37 of the 43 queries",
        "What the approximate index missed amounts to about 1.6 percent",
        "A Wilcoxon test on the queries remaining after ties agrees (p = 0.908)",
    ])
    figure(sl, FIG / "slide2-fig-result-store.png")

    # 10 · Experiment Result — 지표 해설 ① (p.9)
    sl = s[9]; wipe(sl); page_number(sl, 9)
    head(sl, "4. Experiment Result",
         "How the two metrics are read")
    bullets(sl, [
        "Jaccard is how much the two result sets overlap; the p value is how well chance explains it",
    ])
    figure(sl, FIG / "_reading-store.png", top_limit=2.20)

    # 11 · Experiment Result — 위치 요인 (p.10)
    sl = s[10]; wipe(sl); page_number(sl, 10)
    head(sl, "4. Experiment Result",
         "Placement factor — a significant, medium-sized effect")
    bullets(sl, [
        "The same store and the same search method returned mostly different evidence",
        "Since no store effect was detected, the change is attributable to the placement factor",
        "An end-to-end comparison alone would have blamed the store instead",
    ])
    figure(sl, FIG / "slide2-fig-result-placement.png")

    # 12 · Experiment Result — 지표 해설 ② (p.11)
    sl = s[11]; wipe(sl); page_number(sl, 11)
    head(sl, "4. Experiment Result",
         "The same two metrics applied to the placement factor")
    bullets(sl, [
        "Jaccard falls from 0.971 to 0.175 and the p value from 0.523 to 0.003",
    ])
    figure(sl, FIG / "_reading-placement.png", top_limit=2.20)

    # 13 · Experiment Result — 등가성 (p.12)
    sl = s[12]; wipe(sl); page_number(sl, 12)
    head(sl, "4. Experiment Result",
         "Claiming no difference requires an equivalence test of its own")
    bullets(sl, [
        "A non-significant p value is not evidence of equivalence",
        "TOST over a plus or minus 0.15 range gave p = 0.052, so the level was not met",
        "We report no effect detected with a bounded magnitude, not equivalence established",
    ])
    figure(sl, FIG / "_equivalence.png")

    # 14 · Conclusion (p.13)
    sl = s[13]; wipe(sl); page_number(sl, 13)
    head(sl, "5. Conclusion",
         "Moving the placement is a trade in governance rather than in performance")
    y = bullets(sl, [
        "A single bridge condition splits an end-to-end comparison into two single-factor comparisons",
        "No store effect was detected; the entire observed change originated from the placement",
    ], size=16)
    _, tf = textbox(sl, BODY_L, y + 0.22, BODY_W, 0.40)
    put(tf, "A practical decision rule", 17, bold=True, color=BLUE, first=True)
    bullets(sl, [
        "If governance of the stored data is the only goal, replacing the store alone is enough",
        "If the query path must stay inside the boundary, confirm the model size limit first",
        "The bridge condition makes that confirmation possible before the migration",
    ], top=y + 0.72, size=16)

    # 15 · Conclusion — 한계 (p.14)
    sl = s[14]; wipe(sl); page_number(sl, 14)
    head(sl, "5. Conclusion",
         "Limitations and future work")
    y = bullets(sl, [
        "A single Korean corpus of 1,347 chunks drawn from one product domain",
        "Human assessors were not involved in the relevance scoring",
        "Measurements were taken on an entry-level instance, so latency does not generalize",
        "Placement and model size are coupled by construction",
        "The store result does not establish equivalence",
    ], size=16)
    _, tf = textbox(sl, BODY_L, y + 0.22, BODY_W, 0.40)
    put(tf, "Future work", 17, bold=True, color=BLUE, first=True)
    bullets(sl, [
        "Running the smaller model outside the database would separate the two factors",
    ], top=y + 0.72, size=16)

    # 16 · Thank You (p.15)
    sl = s[15]; wipe(sl); page_number(sl, 15)
    # 제목이 없는 장이라 좌상단 강조 사각형만 남으면 제목이 빠진 것처럼 보인다
    for sh in list(sl.shapes):
        if sh.shape_type == 1 and abs(Emu(sh.top).inches - 0.33) < 0.06 \
                and Emu(sh.left).inches < 0.1:
            sh._element.getparent().remove(sh._element)
    _, tf = textbox(sl, 0.74, 2.90, 11.85, 1.60, anchor=MSO_ANCHOR.MIDDLE)
    put(tf, "Thank You", 54, bold=True, color=INK, align=PP_ALIGN.CENTER, first=True)

    prs.save(str(DST))
    return DST


if __name__ == "__main__":
    out = build()
    print(f"저장: {out}")
