"""rev3(사용자 수정본)에 References 슬라이드를 한 장 덧붙인다.

덱 전체를 다시 만들지 않는다. 사용자가 직접 손본 내용을 보존해야 하므로
슬라이드를 새로 하나 만들어 장식만 복제하고 Thank You 앞에 끼워 넣는다.

서지사항은 카메라레디 원고(pygeek2026-en-camera.tex)의 [1]~[9]를 그대로
가져오되, 슬라이드용으로 URL·접속일자만 덜어냈다. 학회지명은 관례대로 축약.
"""
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev3.pptx"
DST = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev4.pptx"

INK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
BODY_L, BODY_W = 0.78, 11.29

REFS = [
    ('[1] P. Lewis et al., “Retrieval-augmented generation for knowledge-intensive '
     'NLP tasks,” Advances in Neural Information Processing Systems, vol. 33, '
     'pp. 9459–9474, 2020.'),
    ('[2] Y. Gao et al., “Retrieval-augmented generation for large language models: '
     'A survey,” arXiv preprint arXiv:2312.10997, 2023.'),
    ('[3] Oracle, “Oracle AI Vector Search User’s Guide,” '
     'Oracle AI Database 26ai documentation, 2026.'),
    ('[4] L. Wang, N. Yang, X. Huang, L. Yang, R. Majumder, and F. Wei, '
     '“Multilingual E5 text embeddings: A technical report,” '
     'arXiv preprint arXiv:2402.05672, 2024.'),
    ('[5] pgvector, “pgvector: Open-source vector similarity search for Postgres,” '
     'GitHub repository, 2026.'),
    ('[6] J. Chen, S. Xiao, P. Zhang, K. Luo, D. Lian, and Z. Liu, '
     '“M3-Embedding: Multi-linguality, multi-functionality, multi-granularity text '
     'embeddings through self-knowledge distillation,” Findings of the Association '
     'for Computational Linguistics: ACL 2024, pp. 2318–2335, 2024.'),
    ('[7] Y. A. Malkov and D. A. Yashunin, “Efficient and robust approximate nearest '
     'neighbor search using hierarchical navigable small world graphs,” IEEE '
     'Transactions on Pattern Analysis and Machine Intelligence, vol. 42, no. 4, '
     'pp. 824–836, 2020.'),
    ('[8] L. Zheng et al., “Judging LLM-as-a-judge with MT-Bench and Chatbot Arena,” '
     'Advances in Neural Information Processing Systems, vol. 36, '
     'Datasets and Benchmarks Track, 2023.'),
    ('[9] D. Lakens, “Equivalence tests: A practical primer for t tests, correlations, '
     'and meta-analyses,” Social Psychological and Personality Science, vol. 8, '
     'no. 4, pp. 355–362, 2017.'),
]


def _style(run, size, bold, color):
    f = run.font
    f.size, f.bold, f.color.rgb, f.name = Pt(size), bold, color, FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "en-US")
    rPr.set("altLang", "en-US")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def build():
    prs = Presentation(str(SRC))
    slides = prs.slides
    donor = slides[13]                      # 본문 슬라이드 하나를 장식 원본으로 쓴다
    thanks = slides[len(slides) - 1]

    new = slides.add_slide(donor.slide_layout)
    for sh in list(new.shapes):             # 레이아웃이 딸려 보낸 빈 placeholder 제거
        sh._element.getparent().remove(sh._element)
    for sh in donor.shapes:                 # 템플릿 장식만 복제
        if sh.shape_type == 1 and sh.name.startswith("직사각형"):
            new.shapes._spTree.append(deepcopy(sh._element))

    # 쪽번호 블록에 새 번호를 직접 써넣는다 (블록 자체가 글자를 품는 구조)
    for sh in new.shapes:
        if abs(Emu(sh.left).inches - 12.66) < .06 and abs(Emu(sh.top).inches - 7.08) < .06:
            tf = sh.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _style(p.add_run(), 12, True, WHITE)
            p.runs[0].text = "15"

    tf = textbox(new, 0.74, 0.30, BODY_W, 0.60)
    _style(tf.paragraphs[0].add_run(), 24, True, INK)
    tf.paragraphs[0].runs[0].text = "References"

    tf = textbox(new, BODY_L, 1.52, BODY_W, 5.20)
    for i, ref in enumerate(REFS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(8)
        p.line_spacing = 1.16
        pPr = p._p.get_or_add_pPr()         # [n] 뒤로 내어쓰기
        pPr.set("marL", str(int(0.42 * 914400)))
        pPr.set("indent", str(int(-0.42 * 914400)))
        _style(p.add_run(), 12, False, INK)
        p.runs[0].text = ref

    # Thank You 앞으로 옮기고, Thank You 쪽번호를 하나 밀어준다
    lst = slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[-1])                     # 방금 만든 References
    lst.insert(len(ids) - 2, ids[-1])
    for sh in thanks.shapes:
        if abs(Emu(sh.left).inches - 12.66) < .06 and abs(Emu(sh.top).inches - 7.08) < .06:
            sh.text_frame.paragraphs[0].runs[0].text = "16"

    prs.save(str(DST))
    return DST


if __name__ == "__main__":
    print(f"저장: {build()}")
