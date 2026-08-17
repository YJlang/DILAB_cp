"""rev9 → rev10. 리허설 피드백 6건 반영 (2026-08-17).

  (1) DILAB 배경을 Introduction 맨 앞으로 — 이전이 왜 필요했는지(운영 환경이
      오라클 표준)를 연구 배경으로 먼저 세운다. 제목·소제목·꼭지도 배경 문맥으로.
  (2) 지표 읽는 법(해설 ①)을 저장소 결과 앞으로 — p값을 쓰기 전에 설명한다.
  (3) 시간·효과 크기 슬라이드 신설 — 8배 저하가 말로만 지나가고 효과 크기는
      설명 없이 그림 구석에만 있었다. 그림은 _latfig.py 로 새로 뽑았다.
  (4) 발표자 노트 18장 전면 교체 — 감사인사·언어 안내 제거, 대명사를 명사로,
      문단을 화면 요소 순서에 맞추고, 암기용 큐(▶) 한 줄을 앞세운다.

슬라이드 레이아웃·영문 본문은 유지한다. 바뀌는 본문은 이동한 DILAB 장의
제목·소제목·둘째 꼭지(배경 문맥화) 뿐이다.
"""
import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev9.pptx"
DST = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev10.pptx"
LATFIG = ROOT / "docs/figures/slide2-fig-latency-effect.png"

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6D, 0x6D, 0x6D)
BLUE = RGBColor(0x33, 0x75, 0xBE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
BODY_L, BODY_W, BODY_T = 0.78, 11.29, 1.52
FIG_BOTTOM = 6.78
SLIDE_W = 13.333

# rev9 인덱스 기준 최종 배열. 17 = 이 스크립트가 맨 끝에 만드는 새 슬라이드.
NEW = 17
FINAL_ORDER = [0, 1, 7, 2, 3, 4, 5, 6, 8, 10, 9, 11, 12, NEW, 13, 14, 15, 16]


def IN(v):
    return Emu(v).inches if v is not None else 0.0


def _style(run, size, bold, color):
    f = run.font
    f.size, f.bold, f.color.rgb, f.name = Pt(size), bold, color, FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "en-US")
    rPr.set("altLang", "en-US")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def bullets(slide, items, top, size=16):
    tf = textbox(slide, BODY_L, top, BODY_W, 0.40 * len(items) + 0.30)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(9)
        p.line_spacing = 1.25
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(0.28 * 914400)))
        pPr.set("indent", str(int(-0.28 * 914400)))
        _style(p.add_run(), size, False, INK)
        p.runs[0].text = "·   " + it
    return top + 0.40 * len(items) + 0.30


def drop(sh):
    sh._element.getparent().remove(sh._element)


def is_decoration(sh):
    return sh.shape_type in (1, 9) and sh.name.startswith(("직사각형", "직선"))


def set_pagenum(slide, n):
    for sh in slide.shapes:
        if (sh.shape_type == 1 and abs(IN(sh.left) - 12.66) < .06
                and abs(IN(sh.top) - 7.08) < .06):
            tf = sh.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _style(p.add_run(), 12, True, WHITE)
            p.runs[0].text = str(n)
            return True
    return False


def box_at(slide, top, tol=.06):
    for sh in slide.shapes:
        if sh.has_text_frame and not is_decoration(sh) and abs(IN(sh.top) - top) < tol:
            return sh
    return None


# ── 발표자 노트 18장 (최종 배열 순서) ─────────────────────────────────────
# 형식: [n/18] 구간 · 제목 / ▶ 암기 큐 / 본문 / → 다음 장.
# 외워서 말하는 전제라, 큐 한 줄만 훑어도 문단 순서가 복원되게 짠다.
NOTES = [
    # 1 · Title
    """[1/18] 0:00–0:25 · Title
▶ 이름 → 제목 → 한 줄 요지
안녕하십니까. 성결대학교 윤준하입니다.
제목은 「벡터 데이터베이스 이전에서, 저장소 교체와 임베딩 위치 효과를 분리하는 요인 분해 방법」입니다.
데이터베이스를 옮겼더니 검색 품질이 달라졌을 때, 무엇 때문에 달라졌는지를 가려내는 방법을 다룹니다.
→ 다음: Contents (0:25)""",
    # 2 · Contents
    """[2/18] 0:25–0:45 · Contents
▶ 배경 → 문제 → 방법 → 실험·결과 → 결론
발표 순서입니다. 이 문제가 나온 서비스 배경을 먼저 보여 드리고, 문제가 무엇인지, 저희 방법이 무엇인지, 실험과 결과, 결론 순서로 진행하겠습니다.
→ 다음: DILAB 배경 (0:45)""",
    # 3 · Introduction — DILAB 배경 (rev9 idx7 이동)
    """[3/18] 0:45–1:30 · Introduction — 실험 무대(DILAB)
▶ 딜랩 소개 → 왼쪽 구조 → 오른쪽 화면 → 왜 이전이 필요했나
연구가 시작된 자리부터 보여 드리겠습니다.
화면은 딜랩이라는 제품 리뷰 분석 데모 서비스입니다. 소비자 리뷰와 전문가 리뷰를 벡터 데이터베이스에 넣어 두고, 질문이 오면 관련 리뷰를 검색해 근거를 인용한 제품 평가 리포트를 만들어 줍니다.
(왼쪽 그림) 왼쪽이 시스템 구조입니다. 리뷰가 임베딩을 거쳐 벡터 저장소에 쌓이고, 질문이 오면 저장소를 검색해 답을 만듭니다.
(오른쪽 그림) 오른쪽이 실제 화면입니다. 다섯 개 축으로 제품을 채점하고, 판단마다 원문 근거가 붙습니다.
그런데 딜랩 데모는 pgvector라는 저장소 위에 만들어져 있고, 실제 운영을 맡을 환경은 오라클 데이터베이스로 표준화되어 있습니다. 실운영으로 가려면 벡터 레이어를 오라클로 옮겨야 했습니다. 오늘 발표할 질문이 바로 그 이전 과정에서 나왔습니다.
→ 다음: 벡터 레이어 (1:30)""",
    # 4 · Introduction — 벡터 레이어
    """[4/18] 1:30–2:10 · Introduction — 벡터 레이어
▶ RAG 세 부분 → 벡터 레이어가 교체 1순위 → 흔한 경로 = 딜랩의 경로
방금 본 딜랩 같은 검색 증강 시스템은 크게 세 부분입니다. 글을 숫자 목록으로 바꾸는 임베딩 모델, 그 숫자를 쌓아 두고 비슷한 것을 찾아 주는 저장소, 찾아온 근거로 답을 쓰는 언어 모델입니다.
(그림 위쪽) 앞의 두 부분, 벡터를 다루는 부분이 서비스를 띄운 뒤 가장 자주 갈아 끼우게 되는 자리입니다.
(그림 아래쪽) 가장 흔한 교체 경로가 화면 아래에 있습니다. 일반 데이터베이스에 벡터 기능을 얹어 쓰다가, 임베딩 계산까지 데이터베이스가 직접 해 주는 제품으로 옮겨 가는 경우입니다. 딜랩의 이전이 바로 이 경로였습니다.
→ 다음: 두 가지가 함께 바뀐다 ★ (2:10)""",
    # 5 · Introduction — 두 가지 ★
    """[5/18] 2:10–3:25 · Introduction — 두 가지가 함께 바뀐다 ★
▶ 바뀌는 게 둘 → ①저장소 ②위치 → 왜 같이 움직이나 → 원인 분리 불가
문제는 이 이전에서 시작됩니다. 한 번 옮길 때 바뀌는 것이 하나가 아니라 둘입니다.
(두 상자) 왼쪽이 pgvector를 얹은 지금 시스템, 오른쪽이 오라클로 옮긴 시스템입니다. 오늘 말씀드릴 내용은 제품의 우열이 아니라, 어느 이전에나 나타나는 구조입니다.
첫째, 저장소가 바뀝니다. 검색 엔진이 바뀌고 찾는 방식도 바뀝니다. 기존 저장소는 빠르게 찾는 대신 가끔 놓치는 방식, 새 저장소는 전부 훑는 방식입니다.
둘째, 임베딩을 계산하는 위치가 바뀝니다. 데이터베이스 바깥에서 하던 계산을 안에서 하게 됩니다.
(붉은 두 줄) 붉게 표시된 두 줄이 실제로 달라지는 항목입니다.
저장소와 위치는 왜 같이 움직일까요. 데이터를 경계 안에 두려고 이전을 결정하는 경우가 많고, 그러면 질문이 밖으로 나가는 것도 막고 싶어져 임베딩이 데이터베이스 안으로 따라 들어옵니다. 그런데 데이터베이스 안에 넣을 수 있는 모델은 크기 제한이 있어서, 대체로 더 작은 모델로 바뀝니다.
결국 이전 전과 후만 비교해서는, 달라진 원인이 저장소인지 모델인지 짚을 수 없습니다.
→ 다음: 원인을 잘못 짚으면 (3:25)""",
    # 6 · Introduction — 원인 잘못
    """[6/18] 3:25–4:05 · Introduction — 원인을 잘못 짚으면
▶ 하락해도 모름 / 그대로여도 모름 → 처방이 어긋남
원인을 못 짚으면 무슨 일이 생길까요.
(왼쪽 상자) 품질이 떨어진 경우, 저장소 탓인지 모델 탓인지 알 수 없습니다.
(오른쪽 상자) 변화가 없는 경우도 마찬가지입니다. 정말 영향이 없었는지, 한쪽이 올리고 한쪽이 내려 상쇄됐는지 구분되지 않습니다.
어느 쪽이든 결론은 "모르겠다"가 됩니다. 그러면 처방이 어긋납니다. 저장소를 범인으로 지목해 되돌려 놓아도 품질은 돌아오지 않습니다. 실제 원인이 임베딩 모델 쪽에 있었다면, 저장소를 되돌리는 일은 헛수고이기 때문입니다.
→ 다음: 중간 조건 ★★ (4:05)""",
    # 7 · Proposed — 브리지 ★★
    """[7/18] 4:05–5:35 · Proposed Method — 중간 조건 ★★
▶ 이사 비유(집=저장소·가구=벡터) → B 만드는 법 → A→B 저장소만 / B→C 위치만 → 재임베딩 불필요
※ 발표의 핵심. 호흡을 두고 그림을 짚으면서 천천히.
저희 제안은 두 조건 사이에 중간 조건 하나를 끼워 넣는 것입니다. 이사에 비유해 보겠습니다.
집을 옮기면서 가구까지 한꺼번에 새로 사면, 새집 분위기가 달라졌을 때 집 때문인지 가구 때문인지 알 수 없습니다. 알아내는 방법은 하나뿐입니다. 먼저 쓰던 가구를 그대로 새집에 옮겨 보고, 그다음에 가구를 바꾸는 것입니다.
이 비유에서 집이 저장소이고, 가구가 벡터입니다.
(가운데 상자) 화면 가운데 조건 B가 "쓰던 가구를 그대로 옮긴 새집"입니다. 지금 시스템이 이미 만들어 둔 벡터를, 다시 계산하지 않고 새 저장소에 그대로 실어 올립니다. 그래서 조건 B는 저장소만 새것이고, 임베딩 모델과 계산 위치는 원래 그대로입니다.
(왼쪽 화살표) A에서 B로 갈 때는 저장소만 달라집니다. (오른쪽 화살표) B에서 C로 갈 때는 임베딩 위치만 달라집니다. 한 번에 하나씩만 보게 됩니다.
(화살표 위 숫자) 화살표 위의 숫자 두 개가 그 두 구간의 결과인데, 잠시 뒤 자세히 말씀드리겠습니다.
실무에서 좋은 점이 하나 있습니다. 벡터를 다시 만들지 않으니 임베딩 비용이 들지 않고, 운영 중인 서비스를 멈추지 않아도 됩니다. 이전을 결정하기 전에 미리 해 볼 수 있다는 점이 이 방법의 쓸모입니다.
→ 다음: 시스템 구조 (5:35)""",
    # 8 · Proposed — 구조
    """[8/18] 5:35–6:15 · Proposed Method — 시스템 구조
▶ 왼쪽 두 번 왕복 → 오른쪽 안에서 끝 → 설정 고정
세 조건이 실제 시스템에서 어떤 모습인지 구조로 보여 드리겠습니다.
(왼쪽) 옮기기 전 시스템은 애플리케이션이 바깥의 임베딩 서비스를 부르고, 받은 벡터를 다시 저장소로 보냅니다. 네트워크 밖으로 두 번 나갔다 옵니다.
(오른쪽) 옮긴 뒤에는 임베딩과 검색이 모두 데이터베이스 안에서 끝납니다. 밖으로 나가는 일은 한 번으로 줄고, 외부 모델 호출은 없어집니다.
꼭 말씀드릴 것이 하나 있습니다. 검색 규칙과 설정값은 세 조건에서 완전히 똑같이 고정했습니다. 비교하려는 요인 말고 다른 것이 섞이면 안 되기 때문입니다.
→ 다음: 실험 설계 (6:15)""",
    # 9 · PerfEval — 실험 설계
    """[9/18] 6:15–7:05 · Performance Evaluation — 실험 설계
▶ 1,347 청크 · 43 질문 → 출처 13+30 → 지표 셋 → 표의 색
※ 이 장을 넘길 때 7분 5초면 정상 속도.
실험 설정입니다.
대상은 앞서 소개드린 딜랩의 리뷰 청크 1,347개이고, 미리 정해 둔 질문 43개를 세 조건에 똑같이 던졌습니다.
질문의 출처를 분명히 말씀드리면, 13개는 실제 서비스 로그에서 뽑았고, 30개는 평가 항목과 구매 단계를 고르게 덮도록 저희가 만들었습니다. 두 묶음이 다르게 행동하지 않는지는 따로 확인했습니다.
지표는 세 가지입니다. 검색 결과가 얼마나 같은지, 찾아온 근거가 질문에 얼마나 맞는지를 0점에서 2점으로 매긴 값, 그리고 검색에 걸린 시간입니다.
(표) 파란 칸이 저장소 요인으로 바뀌는 항목, 붉은 칸이 위치 요인으로 바뀌는 항목입니다. 색이 없는 칸은 세 조건에서 그대로입니다.
〈생략 가능〉 측정은 가장 낮은 등급 인스턴스에서 했습니다.
→ 다음: 지표 읽는 법 (7:05)""",
    # 10 · Result — 지표 읽는 법 (rev9 idx10 이동)
    """[10/18] 7:05–7:55 · Experiment Result — 지표 읽는 법
▶ 겹침 0~1 → p값: 질문 다시 뽑기 반복 → 자주 나오면 우연 / 드물면 진짜
※ 숫자를 처음 다루는 자리. 서두르지 않는다.
결과를 보여 드리기 전에, 결과를 읽을 잣대 두 개를 먼저 설명드리겠습니다.
(왼쪽) 첫째는 두 검색 결과가 얼마나 겹치는지입니다. 상위 10개가 완전히 같으면 1, 하나도 안 겹치면 0입니다.
(오른쪽) 둘째가 p값입니다. 저희가 쓴 질문 43개는 가능한 질문 가운데 일부를 뽑은 것이라, 다른 43개를 뽑았다면 숫자도 조금 달랐을 것입니다. 그래서 두 조건이 사실은 똑같다고 가정해 놓고, 질문 뽑기를 다시 반복하면 지금만 한 차이가 얼마나 자주 나오는지를 봅니다. 그 비율이 p값입니다.
자주 나오는 차이라면 우연으로 설명됩니다. 좀처럼 안 나오는 차이라면 우연이라 보기 어렵습니다.
(그림) 회색 막대들이 우연이 만들 수 있는 차이의 분포이고, 관측된 값이 분포의 어디에 떨어지는지로 판정합니다. 화면의 예시 값은 바로 다음에 보실 저장소 구간입니다.
→ 다음: 저장소 요인 ★ (7:55)""",
    # 11 · Result — 저장소 ★ (rev9 idx9 이동)
    """[11/18] 7:55–9:00 · Experiment Result — 저장소 요인 ★
▶ 결론: 변화 미검출 → 겹침 0.971(43개 중 37개 동일) → p 0.523 → 6개 차이 = 1.6% → 시간은 오히려 단축
첫 번째 결과, 저장소만 바꾼 A에서 B 구간입니다.
결론부터 말씀드리면, 품질 변화가 검출되지 않았습니다.
상위 10개의 겹침이 0.971입니다. 같은 벡터를 실었더니, 질문 43개 가운데 37개에서 상위 10개가 정확히 같았습니다.
근거 품질의 p값은 0.523입니다. 방금 설명드린 잣대로 읽으면, 두 저장소가 똑같다고 해도 절반이 넘는 경우에 나올 만한 차이라는 뜻입니다. 우연으로 충분히 설명됩니다.
나머지 6개 질문에서 차이가 난 이유는, 기존 저장소가 빠르게 찾는 대신 가끔 놓치는 방식이기 때문입니다. 놓친 양은 1.6퍼센트 수준이었습니다.
검색 시간은 150밀리초에서 110밀리초로 오히려 조금 줄었습니다. 저장소만 바꾸는 선택에서는 잃은 것이 없었습니다.
→ 다음: 위치 요인 ★ (9:00)""",
    # 12 · Result — 위치 ★
    """[12/18] 9:00–10:10 · Experiment Result — 위치 요인 ★
▶ 겹침 0.971→0.175 → 품질 1.27→1.01 (p 0.003) → 같은 저장소인데 딴 근거 → 앞뒤만 비교했다면
두 번째 결과, 임베딩 위치만 바꾼 B에서 C 구간입니다. 여기서는 그림이 완전히 달라집니다.
방금 0.971이던 겹침이 0.175로 떨어졌습니다. 거의 다른 목록이 돌아왔다는 뜻입니다.
근거 품질도 2점 만점에 1.27점에서 1.01점으로, 약 20퍼센트 떨어졌습니다. p값은 0.003입니다. 두 조건이 똑같다면 천 번에 세 번 나올 차이라, 우연으로 보기 어렵습니다.
강조하고 싶은 부분은 이것입니다. 저장소도 같고 검색 방식도 같은데, 같은 질문에 대체로 다른 근거가 돌아왔습니다. 앞 구간에서 저장소 효과가 검출되지 않았으므로, 관측된 변화는 사실상 전부 임베딩 위치와, 위치에 딸려 온 모델 크기 제한에서 나온 것입니다.
※ 한 박자 쉬고, 청중을 보며.
만약 이전 전후만 비교했다면 어땠을까요. 저장소 탓으로 돌리고, 저장소를 되돌리고, 품질은 돌아오지 않았을 것입니다.
→ 다음: 지표 해설 ② (10:10)""",
    # 13 · Result — 해설 ②
    """[13/18] 10:10–10:40 · Experiment Result — 지표 해설 ②
▶ 같은 잣대 → 겹침 3/17 → 분포 꼬리 바깥
저장소 구간에 댔던 잣대 두 개를 위치 구간에 그대로 댄 그림입니다.
(왼쪽) 겹침은 0.971에서 0.175로 떨어졌습니다. 두 목록을 합치면 17개인데, 겹치는 것이 3개뿐입니다.
(오른쪽) p값은 0.523에서 0.003으로 내려갔습니다. 관측된 값이, 우연이 만드는 분포의 꼬리 바깥에 나가 있습니다.
같은 잣대를 댔는데 이만큼 다른 그림이 나온다는 것은, 두 요인이 실제로 분리되었다는 뜻입니다.
→ 다음: 시간·효과 크기 (10:40)""",
    # 14 · Result — 시간·효과 크기 (신규)
    """[14/18] 10:40–11:30 · Experiment Result — 시간·효과 크기
▶ 시간 150→110→947(8배) → p는 "우연인가", 효과 크기는 "얼마나 큰가" → 0.098 무시 / 0.473 중간
품질 말고 시간도 보겠습니다.
(왼쪽 패널) 저장소만 바꾼 구간에서는 검색이 150밀리초에서 110밀리초로 오히려 빨라졌습니다. 그런데 임베딩을 안으로 들인 구간에서는 110밀리초에서 947밀리초로, 약 여덟 배 느려졌습니다. 질문이 올 때마다 데이터베이스 안에서 임베딩 모델을 돌려야 하기 때문입니다. 다만 가장 낮은 등급 인스턴스에서 잰 값이라, 장비를 올리면 절대치는 달라질 수 있습니다.
(오른쪽 패널) 품질 차이는 효과 크기라는 값으로도 보겠습니다. p값이 "이 차이가 우연인가"를 묻는 값이라면, 효과 크기는 "그 차이가 실제로 얼마나 큰가"를 묻는 값입니다. 표본이 아주 많으면 하찮은 차이도 p값이 작게 나올 수 있어서, 두 값을 같이 봐야 합니다.
저장소 요인은 0.098로 무시할 수준이고, 위치 요인은 0.473으로 통계에서 중간이라 부르는 크기입니다.
느려진 것도, 품질이 떨어진 것도 모두 위치 요인의 몫이었습니다.
→ 다음: 등가성 검정 (11:30)""",
    # 15 · Result — 등가성
    """[15/18] 11:30–12:30 · Experiment Result — 등가성 검정
▶ 못 찾았다 ≠ 없다 → ±0.15 띠 → 왼끝 이탈, p 0.052 → "나빠졌다 볼 근거 없다"까지만
조심스럽게 말씀드릴 부분이 있습니다.
저장소 요인은 차이가 검출되지 않았다고 했는데, "차이를 못 찾았다"와 "차이가 없다"는 다른 말입니다. 데이터가 부족해서 못 찾았을 수도 있습니다.
"차이가 없다"를 주장하려면 검정을 따로 해야 합니다. 2점 만점에서 플러스마이너스 0.15점 안쪽이면 사실상 같다고 보자는 기준을 정하고, 등가성 검정을 했습니다.
(파란 띠) 옅은 띠가 그 기준입니다. 신뢰구간이 띠 안에 통째로 들어와야 하는데, 왼쪽 끝이 조금 나가 있습니다. p값 0.052로, 기준선 0.05를 아슬아슬하게 넘지 못했습니다.
그래서 결론은 이렇게 보고합니다. "같다는 것이 입증되었다"가 아니라, "차이를 찾지 못했고, 있더라도 이 정도 안쪽이다"입니다.
다만 신뢰구간이 뻗은 방향은 새 저장소 점수가 더 높은 쪽입니다. 저장소를 바꿔서 품질이 나빠졌다고 볼 근거는 없습니다.
→ 다음: Conclusion ★ (12:30)""",
    # 16 · Conclusion ★
    """[16/18] 12:30–13:55 · Conclusion ★
▶ 방법 정리 → 관측 정리 → 거래 문장 → 판단 기준 둘 → 한계
※ 정리와 판단 기준 사이에 호흡 한 번.
정리하겠습니다.
방법 면에서는, 중간에 조건 하나를 두면 앞뒤 비교 한 번이 요인별 비교 두 번으로 나뉩니다.
관측 면에서는, 저장소 교체만으로는 품질 변화가 검출되지 않았고, 관측된 변화는 전부 임베딩 위치에서 나왔습니다.
드리고 싶은 문장은 이것입니다. 임베딩을 데이터베이스 안으로 들이는 것은 성능을 사는 거래가 아니라, 데이터를 경계 안에 두는 대가로 성능을 내주는 거래입니다. 이번 실험에서 치른 값은 근거 품질 20퍼센트와 검색 시간 여덟 배였습니다.
※ 한 박자.
실무 판단 기준으로 바꾸면 이렇습니다. 저장해 둔 데이터만 경계 안에 두면 되는 경우라면, 저장소 교체로 충분합니다. 질의 경로까지 안에 두어야 한다면, 옮기기 전에 모델 크기가 충분한지 확인해야 합니다. 중간 조건은 그 확인을 이전 결정 전에 할 수 있게 해 줍니다.
한계도 분명히 말씀드리겠습니다. 한 도메인의 한국어 데이터 하나에서 얻은 결과이고, 채점에 사람이 참여하지 않았습니다. 시간 수치는 최저 등급 인스턴스 기준이라 일반화할 수 없습니다. 설계상 위치와 모델 크기가 붙어 있어, 둘을 떼어 놓는 네 번째 조건이 다음 과제입니다.
→ 다음: References (13:55)""",
    # 17 · References
    """[17/18] 13:55–14:05 · References
※ 읽지 않는다. 화면만 띄우고 한 문장으로 넘긴다.
참고한 문헌입니다.
→ 다음: Thank You (14:05)""",
    # 18 · Thank You + Q&A 팩
    """[18/18] 14:05–14:20 · Thank You
이상으로 발표를 마치겠습니다. 감사합니다.
〈시간이 남을 때만〉 화면에 메일 주소를 띄워 두었습니다. 편하게 연락 주십시오.
→ 발표 종료 · Q&A 4분
──────────────────────────────
■ 응답 원칙
짧게 시작하고 필요하면 덧붙인다. 모르는 것은 모른다고 말한다.
막히면 이 세 숫자로 돌아온다 — 저장소 0.971 / 위치 0.175 / 시간 8배
──────────────────────────────
■ 예상 질문 · 핵심 한 줄 (전체 답변은 A4 대본 Q1–Q19)
Q1 왜 더 큰 임베딩 모델을 안 썼나 → DB 내장은 모델 크기 제한이 있어 384차원. 그 제약 자체가 이 논문의 대상.
Q2 43개는 적지 않나 → 인정. 그래서 값 하나가 아니라 신뢰구간·효과크기를 함께 보고.
Q3 질의 30개를 직접 만들었는데 믿을 수 있나 → 절대 점수는 실제 로그 쪽이 유의하게 높음(우리가 만든 질의가 더 어려웠음). 단 귀속 근거인 랭킹 일치도는 두 묶음 차이 없음 → 결론 동일.
Q4 사람이 채점 안 했는데 → 한계 인정. 다만 같은 잣대를 세 조건에 동일 적용 → 조건 간 비교 목적에선 방향성 유지.
Q5 새로운 방법인가 → 요소는 익숙함. 기여는 이전 의사결정 상황에서의 절차화 + 재임베딩 없이 무중단으로 가능함을 실측한 점.
Q6 이전하지 말라는 뜻인가 → 아님. 저장 데이터만이면 저장소 교체로 충분, 질의 경로까지면 모델 크기를 먼저 확인.
Q7 저장소 요인에 인덱스 변화가 섞였다 → 인정. recall 98.4% → 인덱스 몫의 상한은 1.6%. 포함해도 결론 유지.
Q8 규모가 커져도 같은가 → 품질은 방향 유지 예상, 속도는 분명히 달라짐. 제안하는 것은 숫자가 아니라 절차.
Q9 8배는 in-DB가 원래 느리다는 뜻인가 → 최저 등급 인스턴스 결과. 임베딩 시간은 A·B 공통으로 1회 측정 → A·B 차이는 검색 시간만 반영.
Q10 검증 방법 → 대응표본 t 주검정 + 윌콕슨 + 부트스트랩 1만, seed 고정. 자카드는 검정이 아니라 요약값.
Q11 다중비교 보정 → 미적용. 핵심 두 개가 0.523/0.003이라 무관. (물으면) 13건 본페로니 0.044.
Q12 ±0.15 마진 근거 → 이론값 아닌 임계값(7.5%). 등가성을 주장하지 않고 오히려 미충족으로 보고.
Q13 0.052면 같다고 봐도 되나 → "같다"는 주장 안 함. "나빠졌다고 볼 근거는 없다"까지.
Q14 왜 0–2 세 단계인가 → LLM 채점 재현성. 필요한 것은 절대값이 아니라 조건 간 차이.
Q15 215건인데 왜 43개로 검정 → 같은 질의의 5점은 독립이 아님 → 질의별 평균 43 대응쌍.
Q16 무엇의 차이인가 → 관련성 평균. 1.228 대 1.270, 차이 0.042.
Q17 우연이란 무엇인가 → 질의 43개를 뽑은 표본의 우연. 채점은 temperature 0이라 무작위 아님.
Q18 왜 이 두 제품인가 → 실제 운영 조합이자 실제 검토 대상. 제품 우열은 판단하지 않음.
Q19 효과 크기란 → 차이의 크기를 표준화한 값. p는 "우연인가", 효과 크기는 "얼마나 큰가". 저장소 0.098 무시 가능 / 위치 0.473 중간.
──────────────────────────────
■ 영어 질문이 들어왔을 때
① 못 알아들었을 때 — 정중한 사과 + 재요청 (이 한 줄만 외워도 됨)
"I'm sorry, my English is not very good. Could you please say that again, a little more slowly?"
(아임 쏘리, 마이 잉글리시 이즈 낫 베리 굿. 쿠쥬 플리즈 쎄이 댓 어겐, 어 리틀 모어 슬로울리?)
"I apologize — I didn't quite catch that. Could you rephrase your question, please?"
"Sorry, one more time, please? I want to make sure I understand your question correctly."
② 반쯤 알아들었을 때 — 확인하고 되묻기 (가장 안전한 수)
"If I understood correctly, you're asking about ___. Is that right?"
"So your question is about the store, not the embedding placement — is that correct?"
→ 상대가 고쳐 주면 그때 답한다. 잘못 알아듣고 엉뚱한 답을 하는 것보다 훨씬 낫다.
③ 그래도 안 되면 — 좌장·통역에게 넘기기 (부끄러운 일이 아님)
"Could I answer this in Korean? I'd like to give you an accurate answer."
"Chair, could I ask for help with the question, please?"
(좌장께 한국어로) "죄송합니다. 질문을 정확히 이해하지 못했는데, 한 번 정리해 주실 수 있을까요?"
④ 답할 수 없는 질문 — 정중한 회피 + 연결
"That's a good question. We haven't measured that directly, so I don't want to guess."
"That goes beyond what our data can support. What I can say is ___."
"Could we continue this after the session? My email is on the slide."
→ "모릅니다"를 영어로 말하는 것은 감점이 아니다. 추측해서 틀리는 것이 감점이다.
⑤ 자주 나올 질문의 짧은 영어 답 (한 문장씩만)
저장소: "Changing only the store showed no detectable change. Top-10 Jaccard was 0.971, p was 0.523."
위치: "Moving the embedding inside carried the entire effect. Jaccard 0.175, p 0.003, about a twenty percent drop in relevance, and search was about eight times slower."
규모: "Our corpus is 1,347 chunks, so I cannot claim this holds at a much larger scale. What we propose is the procedure, not the numbers."
제품: "We do not claim any product is better. We show a procedure for separating the two factors."
⑥ 마무리는 항상 감사로
"Thank you for your question." / "Thank you, that's a helpful point."
──────────────────────────────
■ 마지막 한 줄
천천히 말해도 된다. 침묵 3초는 길게 느껴지지만 청중에게는 짧다.""",
]


def build():
    prs = Presentation(str(SRC))
    assert len(prs.slides._sldIdLst) == 17, "rev9 는 17장이어야 한다"

    # ── (1) DILAB 장(rev9 idx 7)을 배경 문맥으로 고쳐 쓴다 ────────────
    sl = prs.slides[7]
    title = box_at(sl, 0.246)
    assert title.text_frame.text.strip() == "Performance Evaluation"
    title.text_frame.paragraphs[0].runs[0].text = "Introduction"
    sub = box_at(sl, 0.990)
    assert sub.text_frame.text.strip().startswith("Testbed")
    sub.text_frame.paragraphs[0].runs[0].text = (
        "Background — DILAB, a product-review analysis service")
    bl = box_at(sl, BODY_T)
    assert bl.text_frame.text.lstrip().startswith("·")
    # 코퍼스 규모는 실험 설계 장에 그대로 있으므로, 둘째 꼭지를 이전 동기로 바꾼다.
    bl.text_frame.paragraphs[1].runs[0].text = (
        "·   The prototype runs on pgvector; production runs on Oracle, "
        "so the vector layer must migrate")

    # ── (2) 시간·효과 크기 슬라이드 신설 (맨 끝에 만들고 뒤에서 재배열) ──
    donor = prs.slides[13]                 # 등가성 장 — 장식·제목 스타일의 원본
    sl = prs.slides.add_slide(donor.slide_layout)
    for sh in list(sl.shapes):
        drop(sh)
    for sh in donor.shapes:
        if is_decoration(sh):
            sl.shapes._spTree.append(copy.deepcopy(sh._element))
    tf = textbox(sl, BODY_L, 0.246, BODY_W, 0.60)
    _style(tf.paragraphs[0].add_run(), 24, True, INK)
    tf.paragraphs[0].runs[0].text = "Experiment Result"
    tf = textbox(sl, BODY_L, 0.990, BODY_W, 0.42)
    _style(tf.paragraphs[0].add_run(), 17, True, BLUE)
    tf.paragraphs[0].runs[0].text = (
        "Search latency and effect size point to the same factor")
    y = bullets(sl, [
        "Replacing the store left search a little faster — 150 ms down to 110 ms",
        "Moving the embedding inside made search about eight times slower — 110 ms to 947 ms",
        "Effect size states how large a difference is — negligible for the store, "
        "medium for the placement",
    ], BODY_T)
    iw, ih = Image.open(LATFIG).size
    w = 11.60
    h = w / (iw / ih)
    assert y + 0.10 < FIG_BOTTOM - h, "그림이 꼭지를 침범한다"
    sl.shapes.add_picture(str(LATFIG), Inches((SLIDE_W - w) / 2),
                          Inches(FIG_BOTTOM - h), Inches(w), Inches(h))

    # ── (3) 재배열 + 쪽번호 + 노트 ────────────────────────────────────
    lst = prs.slides._sldIdLst
    ids = list(lst)
    assert len(ids) == 18
    for el in ids:
        lst.remove(el)
    for i in FINAL_ORDER:
        lst.append(ids[i])

    assert len(NOTES) == 18
    for i, sl2 in enumerate(prs.slides):
        if i > 0:
            assert set_pagenum(sl2, i), f"슬라이드 {i}: 쪽번호 블록이 없다"
        sl2.notes_slide.notes_text_frame.text = NOTES[i]

    prs.save(str(DST))
    return DST


if __name__ == "__main__":
    print(f"저장: {build()}")
