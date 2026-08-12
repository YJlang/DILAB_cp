"""rev5 → rev6. 랩장님 피드백 3건 반영.

  (1) 꼭지는 페이지당 2~4개  → 1개짜리 3장(7·10·12) 보강, 5~6개짜리 2장(14·15) 축약
  (2) 구조도 페이지 꼭지 1개 → 3개
  (3) 구조도 그림의 고아 '·' 제거 (A·/C· 에서 A,C 만 지워져 점만 남아 있었다)

덱을 다시 만들지 않는다. 사용자가 절 제목의 번호를 직접 떼는 등 손본 내용이
있으므로, 해당 슬라이드의 꼭지 상자와 그림만 갈아 끼운다.
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[3]
SRC = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev5.pptx"
DST = ROOT / "docs/reference/PyGeek2026_발표_윤준하_rev6.pptx"
ARCH = ROOT / "docs/figures/_arch-nodot.png"

INK = RGBColor(0x1A, 0x1A, 0x1A)
BLUE = RGBColor(0x33, 0x75, 0xBE)
FONT = "Arial"
SLIDE_W = 13.333
BODY_L, BODY_W = 0.78, 11.29
BODY_T, FIG_BOTTOM = 1.52, 6.78
GAP = 0.10                      # 꼭지 아래와 그림 사이 최소 간격


def IN(v):
    return Emu(v).inches if v is not None else 0.0


def _style(run, size, bold, color):
    f = run.font
    f.size, f.bold, f.color.rgb, f.name = Pt(size), bold, color, FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "en-US")
    rPr.set("altLang", "en-US")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def box_h(n):
    """기존 판과 같은 산식 — 2꼭지 1.10in, 3꼭지 1.50in."""
    return 0.40 * n + 0.30


def add_bullets(slide, items, top, size=16):
    shp = slide.shapes.add_textbox(Inches(BODY_L), Inches(top),
                                   Inches(BODY_W), Inches(box_h(len(items))))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(9)
        p.line_spacing = 1.25
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(0.28 * 914400)))
        pPr.set("indent", str(int(-0.28 * 914400)))
        _style(p.add_run(), size, False, INK)
        p.runs[0].text = "·   " + it
    return top + box_h(len(items))


def add_subhead(slide, text, top):
    shp = slide.shapes.add_textbox(Inches(BODY_L), Inches(top),
                                   Inches(BODY_W), Inches(0.40))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _style(tf.paragraphs[0].add_run(), 17, True, BLUE)
    tf.paragraphs[0].runs[0].text = text
    return top + 0.50


def drop(sh):
    sh._element.getparent().remove(sh._element)


def bullet_boxes(slide):
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.shape_type == 1:
            continue
        txt = sh.text_frame.text.strip()
        if txt.startswith("·"):
            out.append(sh)
    return sorted(out, key=lambda s: IN(s.top))


def subhead_box(slide, label):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.shape_type != 1 \
                and sh.text_frame.text.strip() == label:
            return sh
    return None


def picture(slide):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            return sh
    return None


def place_picture(slide, png, width=None, bottom=FIG_BOTTOM, top_limit=None):
    iw, ih = Image.open(png).size
    ar = iw / ih
    w = width if width else BODY_W
    h = w / ar
    if top_limit is not None and bottom - h < top_limit:
        h = bottom - top_limit
        w = h * ar
    left = (SLIDE_W - w) / 2
    slide.shapes.add_picture(str(png), Inches(left), Inches(bottom - h),
                             Inches(w), Inches(h))


def build():
    prs = Presentation(str(SRC))
    s = prs.slides

    # ── s7 · Proposed Method (구조도) ──────────────────────────────────
    sl = s[6]
    old = picture(sl)
    old_w = IN(old.width)
    drop(old)
    for b in bullet_boxes(sl):
        drop(b)
    y = add_bullets(sl, [
        "The baseline crosses two network boundaries to reach an external embedding service",
        "After migration, embedding and search both run inside the database boundary",
        "The retrieval rule and its parameters were held identical across the three conditions",
    ], BODY_T)
    place_picture(sl, ARCH, width=old_w, top_limit=y + GAP)

    # ── s10 · 지표 해설 ① ──────────────────────────────────────────────
    sl = s[9]
    pic = picture(sl)
    pw = IN(pic.width)
    drop(pic)
    for b in bullet_boxes(sl):
        drop(b)
    y = add_bullets(sl, [
        "Jaccard measures how much the two top-10 sets overlap, on a 0 to 1 scale",
        "The p value is how well chance alone explains the observed difference",
        "Here both conditions returned the same list and the difference sits mid-distribution",
    ], BODY_T)
    place_picture(sl, ROOT / "docs/figures/_reading-store.png",
                  width=pw, top_limit=y + GAP)

    # ── s12 · 지표 해설 ② ──────────────────────────────────────────────
    sl = s[11]
    pic = picture(sl)
    pw = IN(pic.width)
    drop(pic)
    for b in bullet_boxes(sl):
        drop(b)
    y = add_bullets(sl, [
        "Jaccard falls from 0.971 to 0.175 and the p value from 0.523 to 0.003",
        "Only 3 items are shared between the two top-10 lists",
        "The observed difference now sits far out in the tail of the chance distribution",
    ], BODY_T)
    place_picture(sl, ROOT / "docs/figures/_reading-placement.png",
                  width=pw, top_limit=y + GAP)

    # ── s14 · Conclusion (5 → 4) ──────────────────────────────────────
    sl = s[13]
    sub = subhead_box(sl, "A practical decision rule")
    for b in bullet_boxes(sl):
        drop(b)
    if sub:
        drop(sub)
    y = add_bullets(sl, [
        "A single bridge condition splits an end-to-end comparison into two single-factor comparisons",
        "No store effect was detected; the entire observed change originated from the placement",
    ], BODY_T)
    y = add_subhead(sl, "A practical decision rule", y + 0.22)
    add_bullets(sl, [
        "If governance of the stored data is the only goal, replacing the store alone is enough",
        "If the query path must stay inside, confirm the model size limit before migrating",
    ], y)

    # ── s15 · Conclusion 한계 (6 → 4) ─────────────────────────────────
    sl = s[14]
    sub = subhead_box(sl, "Future work")
    for b in bullet_boxes(sl):
        drop(b)
    if sub:
        drop(sub)
    y = add_bullets(sl, [
        "A single Korean corpus of 1,347 chunks from one domain, scored with no human assessors",
        "Measurements were taken on an entry-level instance, so latency does not generalize",
        "Placement and model size are coupled, and the store result does not establish equivalence",
    ], BODY_T)
    y = add_subhead(sl, "Future work", y + 0.22)
    add_bullets(sl, [
        "Running the smaller model outside the database would separate the two factors",
    ], y)

    prs.save(str(DST))
    return DST


if __name__ == "__main__":
    print(f"저장: {build()}")
